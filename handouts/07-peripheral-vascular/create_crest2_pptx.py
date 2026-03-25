#!/usr/bin/env python3
"""
Generate PowerPoint presentation for CREST-2 Trial
Author: 謝慕揚 MD, PhD, FESC
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor
from datetime import datetime

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
RED = RgbColor(186, 24, 27)
BLUE = RgbColor(0, 114, 188)
DARK_GRAY = RgbColor(64, 64, 64)
LIGHT_GRAY = RgbColor(240, 240, 240)

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RgbColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER

    # Citation
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "N Engl J Med 2026;394:219-31"
    p3.font.size = Pt(18)
    p3.font.color.rgb = BLUE
    p3.alignment = PP_ALIGN.CENTER

    # Author
    txBox4 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = f"整理：謝慕揚 MD, PhD, FESC    {datetime.now().strftime('%Y-%m-%d')}"
    p4.font.size = Pt(16)
    p4.font.color.rgb = DARK_GRAY
    p4.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, bullet_points, highlight_box=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Content area
    content_top = Inches(1.5)
    content_height = Inches(5.5)

    if highlight_box:
        # Add highlight box at top
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), content_top, Inches(12.3), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(255, 250, 205)
        box.line.color.rgb = RgbColor(255, 165, 0)

        hb_tf = box.text_frame
        hb_tf.word_wrap = True
        hb_p = hb_tf.paragraphs[0]
        hb_p.text = highlight_box
        hb_p.font.size = Pt(18)
        hb_p.font.bold = True
        hb_p.font.color.rgb = DARK_GRAY
        hb_p.alignment = PP_ALIGN.CENTER
        hb_tf.paragraphs[0].space_before = Pt(10)

        content_top = Inches(2.9)

    # Bullet points
    txBox2 = slide.shapes.add_textbox(Inches(0.5), content_top, Inches(12.3), content_height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        p.text = "• " + point
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    return slide

def add_two_column_slide(title, left_title, left_points, right_title, right_points):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Left column header
    left_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.6))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = BLUE
    left_header.line.fill.background()
    lh_tf = left_header.text_frame
    lh_p = lh_tf.paragraphs[0]
    lh_p.text = left_title
    lh_p.font.size = Pt(20)
    lh_p.font.bold = True
    lh_p.font.color.rgb = RgbColor(255, 255, 255)
    lh_p.alignment = PP_ALIGN.CENTER

    # Left content
    txBox_left = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(6), Inches(4.5))
    tf_left = txBox_left.text_frame
    tf_left.word_wrap = True
    for i, point in enumerate(left_points):
        if i == 0:
            p = tf_left.paragraphs[0]
        else:
            p = tf_left.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    # Right column header
    right_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6), Inches(0.6))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = BLUE
    right_header.line.fill.background()
    rh_tf = right_header.text_frame
    rh_p = rh_tf.paragraphs[0]
    rh_p.text = right_title
    rh_p.font.size = Pt(20)
    rh_p.font.bold = True
    rh_p.font.color.rgb = RgbColor(255, 255, 255)
    rh_p.alignment = PP_ALIGN.CENTER

    # Right content
    txBox_right = slide.shapes.add_textbox(Inches(6.8), Inches(2.3), Inches(6), Inches(4.5))
    tf_right = txBox_right.text_frame
    tf_right.word_wrap = True
    for i, point in enumerate(right_points):
        if i == 0:
            p = tf_right.paragraphs[0]
        else:
            p = tf_right.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    return slide

def add_results_slide(title, data_rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Table-like content
    y_pos = Inches(1.6)
    for row in data_rows:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y_pos, Inches(12.3), Inches(0.9))
        if row.get('highlight'):
            box.fill.solid()
            box.fill.fore_color.rgb = RgbColor(255, 250, 205)
            box.line.color.rgb = RgbColor(255, 165, 0)
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_GRAY
            box.line.color.rgb = RgbColor(200, 200, 200)

        box_tf = box.text_frame
        box_tf.word_wrap = True
        box_p = box_tf.paragraphs[0]
        box_p.text = row['text']
        box_p.font.size = Pt(20)
        box_p.font.bold = row.get('bold', False)
        box_p.font.color.rgb = RED if row.get('highlight') else DARK_GRAY
        box_p.alignment = PP_ALIGN.CENTER

        y_pos += Inches(1.0)

    return slide

# ============================================================
# Create slides
# ============================================================

# Slide 1: Title
add_title_slide(
    "CREST-2 Trial",
    "Medical Management and Revascularization for Asymptomatic Carotid Stenosis\n無症狀頸動脈狹窄的內科治療與血管重建術"
)

# Slide 2: Key Findings
add_content_slide(
    "核心發現 Key Findings",
    [
        "CAS + 強化內科治療 vs 單獨內科治療：4年事件率 2.8% vs 6.0%",
        "絕對風險差異 3.2%，P=0.02，NNT=31",
        "CEA + 強化內科治療 vs 單獨內科治療：4年事件率 3.7% vs 5.3%",
        "P=0.24，未達統計顯著差異",
    ],
    highlight_box="CAS 顯著降低中風風險；CEA 效益未獲證實"
)

# Slide 3: Background
add_content_slide(
    "研究背景 Background",
    [
        "高度頸動脈狹窄治療方式在國際間差異極大",
        "美國 75-80% 接受 CAS 或 CEA 的患者為無症狀",
        "1990-2000年代試驗顯示 CEA 優於當時的內科治療",
        "現代內科治療、CAS、CEA 技術均有顯著進步",
        "近期小型試驗（SPACE-2、ECST-2）挑戰了 revascularization 的效益"
    ]
)

# Slide 4: Study Design
add_two_column_slide(
    "研究設計 Study Design",
    "Stenting Trial (n=1245)",
    [
        "Medical therapy alone (n=629)",
        "CAS + Medical therapy (n=616)",
        "追蹤中位數 3.6 年"
    ],
    "Endarterectomy Trial (n=1240)",
    [
        "Medical therapy alone (n=623)",
        "CEA + Medical therapy (n=617)",
        "追蹤中位數 4.0 年"
    ]
)

# Slide 5: Inclusion/Exclusion
add_two_column_slide(
    "收案條件 Eligibility Criteria",
    "納入條件 Inclusion",
    [
        "年齡 ≥35 歲",
        "頸動脈狹窄 ≥70%",
        "Doppler PSV ≥230 cm/sec",
        "180天內無同側 stroke/TIA"
    ],
    "排除條件 Exclusion",
    [
        "既往 disabling stroke",
        "不穩定心絞痛",
        "需抗凝治療的心房顫動"
    ]
)

# Slide 6: Medical Management
add_content_slide(
    "強化內科治療 Intensive Medical Management",
    [
        "血壓目標：收縮壓 <130 mmHg",
        "血脂目標：LDL cholesterol <70 mg/dL",
        "2018年後可使用 PCSK9 inhibitor (alirocumab)",
        "血糖、HbA1c 監測與控制",
        "生活型態：戒菸、體重控制、規律運動",
        "電話健康指導 (health coaching)"
    ]
)

# Slide 7: Intervention Protocols
add_two_column_slide(
    "介入治療方案 Intervention Protocols",
    "CAS Protocol",
    [
        "經股動脈入路",
        "必須使用 embolic protection device",
        "術前：Aspirin 325mg + Clopidogrel 75mg BID",
        "術後30天：DAPT",
        "之後：Aspirin 單藥"
    ],
    "CEA Protocol",
    [
        "89% 全身麻醉",
        "術前 ≥48hr：Aspirin 325mg",
        "術中：Heparin 或 Bivalirudin",
        "術後：Aspirin 70-325mg"
    ]
)

# Slide 8: Primary Endpoint
add_content_slide(
    "主要終點 Primary Endpoint",
    [
        "4年複合終點：",
        "  Periprocedural period (Day 0-44): 任何 stroke 或死亡",
        "  Postprocedural period (Day 44-4年): 同側缺血性中風",
        "Stroke 定義：WHO 標準",
        "持續 ≥24 小時的腦功能障礙"
    ]
)

# Slide 9: Main Results - Stenting
add_results_slide(
    "主要結果：Stenting Trial",
    [
        {'text': 'Medical Therapy Alone: 6.0% (95% CI 3.8-8.3)', 'bold': False},
        {'text': 'CAS + Medical Therapy: 2.8% (95% CI 1.5-4.3)', 'bold': False},
        {'text': '絕對風險差異: 3.2% (95% CI 0.6-5.9)', 'bold': True, 'highlight': True},
        {'text': 'P = 0.02   |   NNT = 31', 'bold': True, 'highlight': True},
    ]
)

# Slide 10: Main Results - CEA
add_results_slide(
    "主要結果：Endarterectomy Trial",
    [
        {'text': 'Medical Therapy Alone: 5.3% (95% CI 3.3-7.4)', 'bold': False},
        {'text': 'CEA + Medical Therapy: 3.7% (95% CI 2.1-5.5)', 'bold': False},
        {'text': '絕對風險差異: 1.6% (95% CI -1.1 to 4.3)', 'bold': True},
        {'text': 'P = 0.24 (未達統計顯著)', 'bold': True, 'highlight': True},
    ]
)

# Slide 11: Periprocedural Events
add_two_column_slide(
    "周術期事件 Periprocedural Events (Day 0-44)",
    "Stenting Trial",
    [
        "Medical: 0/629 (0%)",
        "CAS: 8/616 (1.3%)",
        "7 strokes + 1 death"
    ],
    "Endarterectomy Trial",
    [
        "Medical: 3/623 (0.5%)",
        "CEA: 9/617 (1.5%)",
        "All strokes"
    ]
)

# Slide 12: Postprocedural Stroke
add_two_column_slide(
    "術後同側缺血性中風 Postprocedural Ipsilateral Stroke",
    "Stenting Trial",
    [
        "Medical: 1.7%/年 (28 events)",
        "CAS: 0.4%/年 (7 events)",
        "相對風險 4.07 (95% CI 1.78-9.31)"
    ],
    "Endarterectomy Trial",
    [
        "Medical: 1.3%/年 (23 events)",
        "CEA: 0.5%/年 (10 events)",
        "相對風險 2.38 (95% CI 1.13-5.00)"
    ]
)

# Slide 13: Key Interpretation
add_content_slide(
    "結果解讀 Key Interpretation",
    [
        "CAS 的長期同側中風保護 (0.4% vs 1.7%/年) 足以彌補周術期風險",
        "CEA 雖有數值上的優勢，但周術期事件 (1.5%) 抵消部分效益",
        "強化內科治療使年中風率降至 1.3-1.7%（較歷史資料 2-3% 顯著改善）",
        "兩種 revascularization 的術後 ipsilateral stroke 率相似 (0.4-0.5%/年)"
    ]
)

# Slide 14: Limitations
add_content_slide(
    "研究限制 Limitations",
    [
        "非完全雙盲（僅 stroke adjudicators 盲性）",
        "試驗期間內科治療持續進步（血壓目標、PCSK9i、GLP-1 RA）",
        "僅納入經認證的高品質術者，可能無法反映一般實務",
        "未納入 transcarotid artery revascularization (TCAR)",
        "Tipping-point analysis：僅需改變 3-4 例事件可能影響結果"
    ]
)

# Slide 15: Clinical Recommendations
add_content_slide(
    "臨床建議 Clinical Recommendations",
    [
        "所有 ≥70% 無症狀頸動脈狹窄患者應接受強化內科治療",
        "CAS 可考慮用於具經驗團隊、經充分告知的患者",
        "CEA 仍為不適合 CAS 患者的合理選項",
        "定期 carotid duplex 監測狹窄進展",
        "衛教患者 TIA 症狀，若出現應緊急就醫"
    ],
    highlight_box="血壓 <130 mmHg | LDL <70 mg/dL | 戒菸 | 規律運動"
)

# Slide 16: Conclusions
add_content_slide(
    "結論 Conclusions",
    [
        "CAS + 強化內科治療顯著優於單獨內科治療（P=0.02, NNT=31）",
        "CEA + 強化內科治療未達統計顯著效益（P=0.24）",
        "強化內科治療是所有患者的治療基礎",
        "此結果挑戰傳統以 CEA 為首選的治療觀念",
        "持續長期追蹤研究將提供更多證據"
    ],
    highlight_box="CAS 優於 CEA？需更多證據驗證"
)

# Slide 17: References
add_content_slide(
    "參考文獻 References",
    [
        "Brott TG et al. N Engl J Med 2026;394:219-31",
        "Brott TG et al. N Engl J Med 2010;363:11-23 (CREST)",
        "Halliday A et al. Lancet 2021;398:1065-73 (ACST-2)",
        "Reiff T et al. Lancet Neurol 2022;21:877-88 (SPACE-2)",
        "White CJ et al. J Am Coll Cardiol 2022;80:155-70"
    ]
)

# Save presentation
output_path = r"G:\我的雲端硬碟\004 教學資料 at Hsinchu 院內\Journal reading with claude and latex\NEJM CREST-2 2026.pptx"
prs.save(output_path)
print(f"PowerPoint saved to: {output_path}")
